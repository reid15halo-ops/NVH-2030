"""
Add a "Aktuelles Layout" section to NVH2030_Gesamtkonzept.pptx
Inserts:
  - 1 Section Header slide
  - 1 slide with Bild.png  (overview with green markers)
  - 1 slide with Bild (1).png  (full CAD floor plan)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colors (Autoneum theme) ──
LIME     = RGBColor(0xB8, 0xC4, 0x00)
GRAY66   = RGBColor(0x66, 0x66, 0x66)
GRAY_MED = RGBColor(0x99, 0x99, 0x99)
DARK     = RGBColor(0x2D, 0x2D, 0x2D)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
FONT     = "Arial"

# ── Paths ──
BASE     = r"c:\Users\122798\OneDrive - Autoneum\Organized_Workspace\05_Projects\Factory_Layouts\NVH 2030"
PPTX_IN  = os.path.join(BASE, "04_Praesentationen", "NVH2030_Gesamtkonzept.pptx")
IMG1     = os.path.join(BASE, "02_Layouts", "Bild.png")          # overview / annotated
IMG2     = os.path.join(BASE, "02_Layouts", "Bild (1).png")      # full CAD floor plan
OUTPUT   = PPTX_IN   # overwrite in-place (backup kept via OneDrive version history)

# ── Load existing presentation ──
prs = Presentation(PPTX_IN)

# ── Layout references ──
L_SECTION    = prs.slide_layouts[2]   # Section Header
L_TITLE_ONLY = prs.slide_layouts[7]   # Title Only (white)
L_BLANK      = prs.slide_layouts[9]   # Blank

# ── Helper ──
def set_ph(slide, idx, text, font_size=None, bold=None, color=None, alignment=None):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text
            if font_size:
                p.font.size = Pt(font_size)
            if bold is not None:
                p.font.bold = bold
            if color:
                p.font.color.rgb = color
            if alignment:
                p.alignment = alignment
            p.font.name = FONT
            return ph
    return None

def add_textbox(slide, text, left, top, width, height,
                font_size=11, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = alignment
    return txBox

# ════════════════════════════════════════════════════════
# NEW SLIDE A: Section Header — "Aktuelles Layout"
# ════════════════════════════════════════════════════════
sA = prs.slides.add_slide(L_SECTION)
# idx 0 = section title, idx 1 = subtitle (may vary by template)
for ph in sA.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text_frame.clear()
        p = ph.text_frame.paragraphs[0]
        p.text = "Aktuelles Layout"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = FONT
    elif ph.placeholder_format.idx == 1:
        ph.text_frame.clear()
        p = ph.text_frame.paragraphs[0]
        p.text = "Übersicht des bestehenden Fabrikplans — Halle 1 · Halle 2 · Halle 3"
        p.font.size = Pt(16)
        p.font.name = FONT

# ════════════════════════════════════════════════════════
# NEW SLIDE B: Bild.png — annotated overview
# ════════════════════════════════════════════════════════
sB = prs.slides.add_slide(L_TITLE_ONLY)

# Try to set standard placeholders (idx may differ per template)
set_ph(sB, 12, "Aktuelles Layout", font_size=11, color=GRAY66)
set_ph(sB, 0,  "Fabriküberblick — Highlights & Handlungsbedarf", font_size=22, bold=True)
set_ph(sB, 11, "NVH 2030 | Werk Gundernhausen | Autoneum", font_size=8, color=GRAY_MED)

if os.path.exists(IMG1):
    # Fit image across the content area (below title, above footer)
    # Slide is 13.33" x 7.5" (16:9).  Title ~1.5", footer ~0.5"  → image height ~5.3"
    pic = sB.shapes.add_picture(IMG1, Inches(0.3), Inches(1.55), width=Inches(12.7))
    # If picture is taller than available space, constrain height instead
    if pic.height > Inches(5.5):
        sB.shapes._spTree.remove(pic._element)
        sB.shapes.add_picture(IMG1, Inches(0.3), Inches(1.55), height=Inches(5.5))
else:
    add_textbox(sB, f"[Bild fehlt: {IMG1}]", Inches(0.5), Inches(2), Inches(12), Inches(1),
                color=RGBColor(0xFF, 0, 0))

# ════════════════════════════════════════════════════════
# NEW SLIDE C: Bild (1).png — full CAD floor plan
# ════════════════════════════════════════════════════════
sC = prs.slides.add_slide(L_TITLE_ONLY)

set_ph(sC, 12, "Aktuelles Layout", font_size=11, color=GRAY66)
set_ph(sC, 0,  "Gesamtlayoutplan — Halle 1 · Halle 2 · Halle 3 (CAD-Auszug)", font_size=22, bold=True)
set_ph(sC, 11, "NVH 2030 | Werk Gundernhausen | Autoneum", font_size=8, color=GRAY_MED)

if os.path.exists(IMG2):
    pic2 = sC.shapes.add_picture(IMG2, Inches(0.3), Inches(1.55), width=Inches(12.7))
    if pic2.height > Inches(5.5):
        sC.shapes._spTree.remove(pic2._element)
        sC.shapes.add_picture(IMG2, Inches(0.3), Inches(1.55), height=Inches(5.5))
else:
    add_textbox(sC, f"[Bild fehlt: {IMG2}]", Inches(0.5), Inches(2), Inches(12), Inches(1),
                color=RGBColor(0xFF, 0, 0))

# ── Save ──
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Total slides now: {len(prs.slides)}")
