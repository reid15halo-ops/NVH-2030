"""
Patch slides 17 and 18 of NVH2030_Gesamtkonzept.pptx:
  - Resize images to leave a right-hand text column
  - Add descriptive text / legend beside each image
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
LIME     = RGBColor(0xB8, 0xC4, 0x00)
GRAY66   = RGBColor(0x66, 0x66, 0x66)
GRAY_MED = RGBColor(0x99, 0x99, 0x99)
DARK     = RGBColor(0x2D, 0x2D, 0x2D)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RED_M    = RGBColor(0xCC, 0x00, 0x00)   # red machines
GREEN_M  = RGBColor(0x00, 0xAA, 0x44)   # green machines
ORANGE   = RGBColor(0xE8, 0x6A, 0x10)
BLUE_L   = RGBColor(0x69, 0xBB, 0xFF)
CARD_BG  = RGBColor(0xF0, 0xF0, 0xF0)
FONT     = "Arial"

BASE   = r"c:\Users\122798\OneDrive - Autoneum\Organized_Workspace\05_Projects\Factory_Layouts\NVH 2030"
PPTX   = os.path.join(BASE, "04_Praesentationen", "NVH2030_Gesamtkonzept.pptx")
IMG1   = os.path.join(BASE, "02_Layouts", "Bild.png")
IMG2   = os.path.join(BASE, "02_Layouts", "Bild (1).png")

prs = Presentation(PPTX)

# ── Helpers ──
def add_tb(slide, text, left, top, width, height,
           font_size=11, bold=False, color=DARK,
           alignment=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = alignment
    return tb

def add_bullets(slide, items, left, top, width, height,
                font_size=11, spacing=Pt(5), color=DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (bullet, clr) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.name = FONT
        p.font.color.rgb = clr if clr else color
        p.space_after = spacing
    return tb

def add_card(slide, left, top, width, height, fill=CARD_BG, border=None, accent=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border:
        card.line.color.rgb = border
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    if accent:
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()
    return card

def legend_dot(slide, color, label, left, top, font_size=10):
    """Small colored circle + label."""
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.03),
                                  Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_tb(slide, label, left + Inches(0.26), top, Inches(2.5), Inches(0.28),
           font_size=font_size, color=DARK)

def remove_pictures(slide):
    """Remove all picture shapes from a slide."""
    to_remove = [sh for sh in slide.shapes if sh.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE = 13
    for sh in to_remove:
        sp = sh._element
        sp.getparent().remove(sp)

# ════════════════════════════════════════════════════════════════
# SLIDE 17  (index 17) — Bild.png  — Highlights & Handlungsbedarf
# ════════════════════════════════════════════════════════════════
s17 = prs.slides[17]
remove_pictures(s17)

# Image: left 2/3 of slide (0.3" → 8.7" wide, from y=1.55")
IMG_H = Inches(5.3)
if os.path.exists(IMG1):
    s17.shapes.add_picture(IMG1, Inches(0.3), Inches(1.55), width=Inches(8.6))

# Right column: x=9.1"  width=4.0"
RX = Inches(9.1)
RW = Inches(4.0)

# Section label
add_tb(s17, "Was zeigt diese Ansicht?", RX, Inches(1.6), RW, Inches(0.35),
       font_size=13, bold=True, color=LIME)

# Context paragraphs
context = [
    ("Die grünen Markierungen zeigen Bereiche und Maschinen, die im Rahmen "
     "von NVH 2030 Phase A aus Halle 3 in Halle 2 verlagert werden.",
     DARK),
    ("", None),
    ("Jede Ellipse entspricht einer Maschinengruppe oder einem Produktionsbereich, "
     "der nach der Konsolidierung wegfällt oder zusammengeführt wird.", GRAY66),
]
add_bullets(s17, context, RX, Inches(2.05), RW, Inches(1.8), font_size=10, color=DARK)

# Divider line
line = s17.shapes.add_shape(MSO_SHAPE.RECTANGLE, RX, Inches(3.95), RW, Inches(0.02))
line.fill.solid()
line.fill.fore_color.rgb = GRAY_MED
line.line.fill.background()

# Key facts cards
add_tb(s17, "Kennzahlen Phase A", RX, Inches(4.05), RW, Inches(0.3),
       font_size=11, bold=True, color=DARK)

facts = [
    ("5 Maschinen abzubauen", LIME),
    ("Halle 3 vollständig räumen", LIME),
    ("Max. 2 Wochen Stillstand (Sommerpause)", ORANGE),
    ("Investition: 262 – 460 k €", BLUE_L),
]
y = Inches(4.4)
for txt, clr in facts:
    add_card(s17, RX, y, RW, Inches(0.38), accent=clr)
    add_tb(s17, "  " + txt, RX + Inches(0.15), y + Inches(0.06),
           RW - Inches(0.2), Inches(0.28), font_size=10, color=DARK)
    y += Inches(0.44)

# Footer note
add_tb(s17, "Quelle: CAD-Layoutplan Werk Gundernhausen  |  Stand: Feb 2026",
       RX, Inches(6.85), RW, Inches(0.25), font_size=8, color=GRAY_MED)

# ════════════════════════════════════════════════════════════════
# SLIDE 18  (index 18) — Bild (1).png  — Gesamtlayoutplan CAD
# ════════════════════════════════════════════════════════════════
s18 = prs.slides[18]
remove_pictures(s18)

if os.path.exists(IMG2):
    s18.shapes.add_picture(IMG2, Inches(0.3), Inches(1.55), width=Inches(8.6))

# Right column
add_tb(s18, "Legende", RX, Inches(1.6), RW, Inches(0.35),
       font_size=13, bold=True, color=LIME)

add_tb(s18, "Farbkodierung im CAD-Plan:", RX, Inches(2.0), RW, Inches(0.28),
       font_size=10, bold=True, color=DARK)

legend_dot(s18, RED_M,   "Rot  — Maschinen / Anlagen werden abgebaut",   RX, Inches(2.35), 10)
legend_dot(s18, GREEN_M, "Grün — Maschinen / Anlagen verbleiben oder werden verlagert", RX, Inches(2.7), 10)

# Divider
line2 = s18.shapes.add_shape(MSO_SHAPE.RECTANGLE, RX, Inches(3.1), RW, Inches(0.02))
line2.fill.solid()
line2.fill.fore_color.rgb = GRAY_MED
line2.line.fill.background()

# Hallen overview
add_tb(s18, "Hallen-Übersicht", RX, Inches(3.2), RW, Inches(0.3),
       font_size=11, bold=True, color=DARK)

hallen = [
    ("Halle 1", "Bestehende Produktion bleibt unverändert.\nKeine Eingriffe in Phase A.", GRAY66),
    ("Halle 2", "Zielzustand: alle NVH-Anlagen konsolidiert.\nC-Linie + F-Linie + neue CO₂-Infrastruktur.", LIME),
    ("Halle 3", "Wird vollständig geräumt.\nStrategisch freigezogen für Phase B.", ORANGE),
]

y = Inches(3.6)
for name, desc, clr in hallen:
    add_card(s18, RX, y, RW, Inches(0.85), accent=clr)
    add_tb(s18, name, RX + Inches(0.15), y + Inches(0.06),
           RW - Inches(0.2), Inches(0.28), font_size=11, bold=True, color=clr)
    add_tb(s18, desc, RX + Inches(0.15), y + Inches(0.35),
           RW - Inches(0.2), Inches(0.46), font_size=9, color=GRAY66)
    y += Inches(0.95)

# Footer
add_tb(s18, "Quelle: CAD-Layoutplan Werk Gundernhausen  |  Stand: Feb 2026",
       RX, Inches(6.85), RW, Inches(0.25), font_size=8, color=GRAY_MED)

# ── Save ──
prs.save(PPTX)
print(f"Patched & saved: {PPTX}")
print(f"Total slides: {len(prs.slides)}")
