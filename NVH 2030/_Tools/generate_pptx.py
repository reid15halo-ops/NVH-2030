"""
NVH 2030 — Management Summary PPTX
Uses Autoneum corporate POTX template (200624_ATN_Template_16to9.potx)
Max 10 slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os, copy

# ── Autoneum Theme Colors (from template) ──
LIME      = RGBColor(0xB8, 0xC4, 0x00)  # accent1
GRAY66    = RGBColor(0x66, 0x66, 0x66)  # accent2
BLUE_L    = RGBColor(0x69, 0xBB, 0xFF)  # accent3
ORANGE    = RGBColor(0xE8, 0x6A, 0x10)  # accent4
GRAY_DC   = RGBColor(0xDC, 0xDC, 0xDC)  # accent5/lt2
BLACK     = RGBColor(0x00, 0x00, 0x00)  # dk1
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)  # lt1
DARK      = RGBColor(0x2D, 0x2D, 0x2D)
GRAY_MED  = RGBColor(0x99, 0x99, 0x99)
CARD_BG   = RGBColor(0xF0, 0xF0, 0xF0)  # light card on white

# Font: Arial (from template theme)
FONT = "Arial"

# ── Paths ──
BASE = r"c:\Users\122798\OneDrive - Autoneum\Organized_Workspace\05_Projects\Factory_Layouts\NVH 2030"
TEMPLATE = os.path.join(BASE, "_Tools", "_template_converted.pptx")
IMG_V2 = os.path.join(BASE, "02_Layouts", "Vorschlag 2 24-02-2026.png")
IMG_V3 = os.path.join(BASE, "02_Layouts", "Vorschlag 3 24-02-2026.png")
IMG_V4 = os.path.join(BASE, "02_Layouts", "Vorschlag 4 24-02-2026.jpg")
OUTPUT = os.path.join(BASE, "04_Praesentationen", "NVH2030_Management_Summary.pptx")

# ── Load template ──
prs = Presentation(TEMPLATE)

# ── Delete existing 3 template slides ──
# Just remove the sldId entries from the presentation XML.
# The slide parts remain as orphans but won't appear.
ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
ns_r = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
ct_pres = prs.element  # CT_Presentation
sldIdLst = ct_pres.find(ns_p + 'sldIdLst')
for sldId in list(sldIdLst):
    sldIdLst.remove(sldId)

# ── Layout references ──
L_TITLE      = prs.slide_layouts[0]   # Title Slide (white)
L_TITLE_GREY = prs.slide_layouts[1]   # Title Slide (grey)
L_SECTION    = prs.slide_layouts[2]   # Section Header
L_CONTENT    = prs.slide_layouts[3]   # Title and Content (white)
L_CONTENT_G  = prs.slide_layouts[4]   # Title and Content (grey)
L_TWO_COL    = prs.slide_layouts[5]   # Two Content (white)
L_TWO_COL_G  = prs.slide_layouts[6]   # Two Content (grey)
L_TITLE_ONLY = prs.slide_layouts[7]   # Title Only (white)
L_TITLE_O_G  = prs.slide_layouts[8]   # Title Only (grey)
L_BLANK      = prs.slide_layouts[9]   # Blank
L_CLOSING    = prs.slide_layouts[10]  # Closing Slide

# ── Helper Functions ──

def set_ph(slide, idx, text, font_size=None, bold=None, color=None, alignment=None):
    """Set placeholder text with formatting."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text
            if font_size:
                p.font.size = Pt(font_size)
            if bold is not None:
                p.font.bold = bold
            if color:
                p.font.color.rgb = color
            if alignment:
                p.alignment = alignment
            p.font.name = FONT
            return ph
    return None


def set_ph_bullets(slide, idx, items, font_size=11, color=None, bold_first=False, spacing=Pt(6)):
    """Set placeholder with multiple bullet paragraphs."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = item
                p.font.size = Pt(font_size)
                p.font.name = FONT
                p.space_after = spacing
                if color:
                    p.font.color.rgb = color
                if bold_first and i == 0:
                    p.font.bold = True
            return ph
    return None


def add_textbox(slide, text, left, top, width, height, font_size=11, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, name=FONT):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = alignment
    return txBox


def add_bullet_textbox(slide, items, left, top, width, height, font_size=11, color=BLACK, spacing=Pt(4)):
    """Add a textbox with bullet items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # Use bullet character
        p.text = "\u2022  " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = spacing
    return txBox


def add_card(slide, left, top, width, height, fill_color=CARD_BG, border_color=None, accent_top=None):
    """Add a rounded rectangle card."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    if accent_top:
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.05))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent_top
        stripe.line.fill.background()
    return card


# ════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(L_TITLE)
set_ph(s1, 0, "NVH 2030 \u2014 Zukunftssicheres Fabriklayout", font_size=32, bold=True)
set_ph(s1, 1, "Management Summary  |  Werk Gundernhausen  |  26.02.2026", font_size=14)

# ════════════════════════════════════════════════════════
# SLIDE 2: Ausgangslage
# ════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(L_CONTENT)
set_ph(s2, 12, "Ausgangslage", font_size=11, color=GRAY66)
set_ph(s2, 0, "Warum muss sich das Layout andern?", font_size=24, bold=True)
set_ph(s2, 11, "NVH 2030 | Werk Gundernhausen | Autoneum", font_size=8, color=GRAY_MED)

items = [
    "NVH-Geschaft schrumpft \u2014 Maschinenauslastung sinkt kontinuierlich",
    "Produktion aktuell uber 2 Hallen verteilt (Halle 2 + Halle 3)",
    "Ineffiziente Transportwege und Zwischenlagerung von Halbteilen zwischen Hallen",
    "Steigende Energiekosten fur Formtragerheizung in 2 getrennten Hallen",
    "Halle 3 bindet wertvolle Flache ohne strategischen Nutzen",
    "",
    "Eckdaten:",
    "   Jahresvolumen: 500\u20131.000 t   |   Mitarbeiter: < 30   |   5 Maschinen zu entfernen",
]
set_ph_bullets(s2, 1, items, font_size=13, spacing=Pt(6))

# ════════════════════════════════════════════════════════
# SLIDE 3: Three Layout Options with images
# ════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(L_TITLE_ONLY)
set_ph(s3, 12, "Vergleich", font_size=11, color=GRAY66)
set_ph(s3, 0, "Drei Layoutvorschlage", font_size=24, bold=True)
set_ph(s3, 11, "NVH 2030 | Werk Gundernhausen | Autoneum", font_size=8, color=GRAY_MED)

# Vorschlag 2
add_textbox(s3, "Vorschlag 2", Inches(0.4), Inches(1.6), Inches(3.8), Inches(0.3),
            font_size=12, bold=True, color=GRAY66)
if os.path.exists(IMG_V2):
    s3.shapes.add_picture(IMG_V2, Inches(0.4), Inches(1.95), width=Inches(3.6))

# Vorschlag 3
add_textbox(s3, "Vorschlag 3", Inches(4.5), Inches(1.6), Inches(3.8), Inches(0.3),
            font_size=12, bold=True, color=GRAY66)
if os.path.exists(IMG_V3):
    s3.shapes.add_picture(IMG_V3, Inches(4.5), Inches(1.95), width=Inches(3.6))

# Vorschlag 4 — highlighted
# Lime border
highlight = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.45), Inches(1.48), Inches(4.4), Inches(5.5))
highlight.fill.background()
highlight.line.color.rgb = LIME
highlight.line.width = Pt(3)

add_textbox(s3, "Vorschlag 4  \u2014  AUSGEWAHLT", Inches(8.6), Inches(1.6), Inches(4), Inches(0.3),
            font_size=12, bold=True, color=LIME)
if os.path.exists(IMG_V4):
    s3.shapes.add_picture(IMG_V4, Inches(8.6), Inches(1.95), width=Inches(4.0))

# ════════════════════════════════════════════════════════
# SLIDE 4: Vorschlag 4 Detail
# ════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(L_TITLE_ONLY)
set_ph(s4, 12, "Entscheidung", font_size=11, color=GRAY66)
set_ph(s4, 0, "Vorschlag 4 \u2014 Konsolidierung in Halle 2", font_size=24, bold=True)
set_ph(s4, 11, "NVH 2030 | Werk Gundernhausen | Autoneum", font_size=8, color=GRAY_MED)

# Large layout image left
if os.path.exists(IMG_V4):
    s4.shapes.add_picture(IMG_V4, Inches(0.4), Inches(1.6), width=Inches(7.0))

# Zone cards right
zones = [
    ("Produktbereich 1", "Gemischte Produktionszellen\nProduktgruppe A\nHL spezifisch zugeordnet", LIME),
    ("Produktbereich 2 \u2014 Hauptproduktion", "Schaumlinien C + F (je 6 Mischkopfe)\nCO2 + Trennmittel-Bereiche\nDosier- & Mischanlagen", BLUE_L),
    ("Produktbereich 3", "Gemischte Produktionszellen\nProduktgruppe C\nAnbindung Versand", ORANGE),
]

y = Inches(1.7)
for title, desc, color in zones:
    add_card(s4, Inches(7.8), y, Inches(5.0), Inches(1.35), fill_color=CARD_BG, accent_top=color)
    add_textbox(s4, title, Inches(8.0), y + Inches(0.12), Inches(4.6), Inches(0.3),
                font_size=12, bold=True, color=color)
    add_textbox(s4, desc, Inches(8.0), y + Inches(0.45), Inches(4.6), Inches(0.85),
                font_size=10, color=GRAY66)
    y += Inches(1.55)

# Staplerweg note
add_textbox(s4, "Staplerverkehrswege trennen alle drei Bereiche sicher voneinander",
            Inches(7.8), y + Inches(0.1), Inches(5.0), Inches(0.3),
            font_size=10, color=GRAY_MED)

# ════════════════════════════════════════════════════════
# SLIDE 5: Why Vorschlag 4?
# ════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(L_TITLE_ONLY)
set_ph(s5, 12, "Argumentation", font_size=11, color=GRAY66)
set_ph(s5, 0, "Warum Vorschlag 4?  \u2014  12 Argumente", font_size=24, bold=True)
set_ph(s5, 11, "NVH 2030 | Werk Gundernhausen | Autoneum", font_size=8, color=GRAY_MED)

categories = [
    ("LOGISTIK", LIME, [
        "Kurzere Transportwege",
        "Kein Zwischenlager Halbteile",
        "Eindeutige HL-Zuordnung",
        "Platz fur Container + Werkzeuge"
    ]),
    ("EFFIZIENZ", BLUE_L, [
        "One-Way-Flow \u2014 weniger Personal",
        "F-Linie entlastet C-Linie",
        "C + F je 6 Mischkopfe",
        "CO2 bedient beide Anlagen"
    ]),
    ("ENERGIE", ORANGE, [
        "Formtragerheizung nur 1 Halle",
        "Halle 3 komplett abschaltbar",
    ]),
    ("STRATEGIE", GRAY66, [
        "Halle 3 frei fur Zukunft",
        "CO2-RA Optionalitat (Phase B)",
    ]),
]

x = Inches(0.4)
col_w = Inches(3.0)
for cat_name, color, items in categories:
    # Category header bar
    hdr = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.6), col_w, Inches(0.35))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    add_textbox(s5, cat_name, x + Inches(0.08), Inches(1.62), col_w - Inches(0.16), Inches(0.3),
                font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Items
    y = Inches(2.1)
    for item in items:
        add_card(s5, x, y, col_w, Inches(0.42), fill_color=CARD_BG)
        add_textbox(s5, "\u2022  " + item, x + Inches(0.1), y + Inches(0.05), col_w - Inches(0.2), Inches(0.35),
                    font_size=10, color=DARK)
        y += Inches(0.48)

    x += Inches(3.15)

# ════════════════════════════════════════════════════════
# SLIDE 6: Cost Estimate Phase A
# ════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(L_TITLE_ONLY)
set_ph(s6, 12, "Budget", font_size=11, color=GRAY66)
set_ph(s6, 0, "Kostenschatzung Phase A", font_size=24, bold=True)
set_ph(s6, 11, "NVH 2030 | Werk Gundernhausen | Autoneum", font_size=8, color=GRAY_MED)

cost_items = [
    ("A1", "Maschinenabbau\n& Halle 3 Raumung", "81\u2013132k \u20ac", LIME),
    ("A2", "Anlagenumstellung\nHalle 2", "58\u2013100k \u20ac", BLUE_L),
    ("A3", "Allgemeine\nInfrastruktur", "115\u2013210k \u20ac", ORANGE),
    ("A4", "Wassrige\nTrennmittel", "8\u201318k \u20ac", GRAY66),
]

x = Inches(0.4)
for code, label, cost, color in cost_items:
    add_card(s6, x, Inches(1.7), Inches(2.9), Inches(2.2), fill_color=CARD_BG, accent_top=color)
    add_textbox(s6, code, x + Inches(0.15), Inches(1.85), Inches(2.6), Inches(0.25),
                font_size=11, bold=True, color=color)
    add_textbox(s6, label, x + Inches(0.15), Inches(2.15), Inches(2.6), Inches(0.6),
                font_size=10, color=GRAY66)
    add_textbox(s6, cost, x + Inches(0.15), Inches(2.95), Inches(2.6), Inches(0.4),
                font_size=18, bold=True, color=DARK)
    x += Inches(3.1)

# Total bar
total_bar = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(4.3), Inches(8.0), Inches(1.0))
total_bar.fill.solid()
total_bar.fill.fore_color.rgb = LIME
total_bar.line.fill.background()

add_textbox(s6, "GESAMT PHASE A", Inches(2.8), Inches(4.4), Inches(3), Inches(0.3),
            font_size=12, bold=True, color=DARK)
add_textbox(s6, "262.000 \u2013 460.000 \u20ac", Inches(2.8), Inches(4.7), Inches(7.5), Inches(0.5),
            font_size=26, bold=True, color=DARK)

# Notes
notes = [
    "Sommerpause nutzbar (max. 2 Wochen Stillstand)",
    "Keine Genehmigungen erforderlich",
    "Personalkosten und Produktionsausfallkosten nicht enthalten",
]
add_bullet_textbox(s6, notes, Inches(0.4), Inches(5.6), Inches(12), Inches(1.3),
                   font_size=10, color=GRAY66)

# ════════════════════════════════════════════════════════
# SLIDE 7: Roadmap
# ════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(L_TITLE_ONLY)
set_ph(s7, 12, "Umsetzungsfahrplan", font_size=11, color=GRAY66)
set_ph(s7, 0, "Roadmap Phase A", font_size=24, bold=True)
set_ph(s7, 11, "NVH 2030 | Werk Gundernhausen | Autoneum", font_size=8, color=GRAY_MED)

# Timeline bar
timeline_y = Inches(2.4)
tl = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), timeline_y, Inches(11.5), Inches(0.06))
tl.fill.solid()
tl.fill.fore_color.rgb = GRAY_DC
tl.line.fill.background()

milestones = [
    ("Q1 2026", "Layout bestatigen\nBudgetfreigabe", Inches(1.0), LIME),
    ("Q2 2026", "Detailplanung\nAusschreibungen", Inches(3.6), BLUE_L),
    ("Q3 2026", "Sommerpause\nKritische Umbauten", Inches(6.2), ORANGE),
    ("Q4 2026", "Maschinenabbau\nHalle 3 raumen", Inches(8.8), GRAY66),
    ("Q1 2027", "Abschluss\nPhase A fertig", Inches(11.0), LIME),
]

for label, desc, x, color in milestones:
    # Dot
    dot = s7.shapes.add_shape(MSO_SHAPE.OVAL, x, timeline_y - Inches(0.1), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    # Quarter above
    add_textbox(s7, label, x - Inches(0.35), timeline_y - Inches(0.55), Inches(1.2), Inches(0.3),
                font_size=12, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # Description below
    add_textbox(s7, desc, x - Inches(0.5), timeline_y + Inches(0.3), Inches(1.8), Inches(0.7),
                font_size=10, color=GRAY66, alignment=PP_ALIGN.CENTER)

# Key constraints below
constraints = [
    "Max. 2 Wochen Produktionsstillstand (Sommerpause 2026)",
    "Keine behordlichen Genehmigungen erforderlich",
    "5 Maschinen abzubauen",
    "Phase B (CO2-RA, Zukunftsoptionen) erst nach Abschluss Phase A",
]
add_bullet_textbox(s7, constraints, Inches(1.5), Inches(4.2), Inches(10), Inches(2.5),
                   font_size=12, color=DARK, spacing=Pt(8))

# ════════════════════════════════════════════════════════
# SLIDE 8: Phase B — Future Options
# ════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(L_TITLE_ONLY)
set_ph(s8, 12, "Zukunft", font_size=11, color=GRAY66)
set_ph(s8, 0, "Halle 3 \u2014 Zukunftsoptionen (Phase B)", font_size=24, bold=True)
set_ph(s8, 11, "NVH 2030 | Werk Gundernhausen | Autoneum", font_size=8, color=GRAY_MED)

# Disclaimer bar
disc = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.35))
disc.fill.solid()
disc.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xCD)
disc.line.fill.background()
add_textbox(s8, "Alle Optionen sind Explorationen aus der R&D-Pipeline. Keine ist fur Gundernhausen bestatigt oder budgetiert.",
            Inches(0.6), Inches(1.52), Inches(12), Inches(0.3),
            font_size=10, bold=True, color=ORANGE)

# 4 Options
options = [
    ("1", "CO2-Trennmittel", "TRL 6-7", "230\u2013365k \u20ac", "Ellzee validiert\nCO2-Tank vorhanden\n~4,4 J. Amortisation", LIME),
    ("2", "PU-Chem. Recycling", "TRL 3-4", "500k\u20131M CHF", "Selective Polymer-Retention\n100% Baumwolle + 97% PET\nEU Green Deal", BLUE_L),
    ("3", "LHS-TPC Leichtbau", "TRL 4-5", "Offen", "PP-GF Hohlstruktur-Panels\n100% recycelbar\nErsatz IboComb-Set", ORANGE),
    ("4", "FR-PU Batterieschutz", "TRL 4-5", "Offen", "In-Press FR-PU\nErsatz Mica-Shields\nBEV-Markt wachsend", GRAY66),
]

x = Inches(0.4)
for num, title, trl, invest, desc, color in options:
    add_card(s8, x, Inches(2.1), Inches(3.0), Inches(4.6), fill_color=CARD_BG, accent_top=color)

    # Number circle
    circle = s8.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.12), Inches(2.25), Inches(0.35), Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_textbox(s8, num, x + Inches(0.12), Inches(2.26), Inches(0.35), Inches(0.35),
                font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(s8, title, x + Inches(0.55), Inches(2.25), Inches(2.3), Inches(0.3),
                font_size=12, bold=True, color=DARK)

    # TRL + Investment
    add_textbox(s8, trl + "   |   Invest: " + invest, x + Inches(0.12), Inches(2.7), Inches(2.7), Inches(0.25),
                font_size=9, color=color)

    # Description
    add_textbox(s8, desc, x + Inches(0.12), Inches(3.05), Inches(2.7), Inches(3.5),
                font_size=10, color=GRAY66)

    x += Inches(3.15)

# ════════════════════════════════════════════════════════
# SLIDE 9: Strategic Assessment
# ════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(L_TWO_COL)
set_ph(s9, 12, "Bewertung", font_size=11, color=GRAY66)
set_ph(s9, 0, "Strategische Bewertung", font_size=24, bold=True)
set_ph(s9, 11, "NVH 2030 | Werk Gundernhausen | Autoneum", font_size=8, color=GRAY_MED)

# Left: Phase A
items_a = [
    "PHASE A \u2014 KERNPROJEKT",
    "",
    "Investition rechnet sich allein durch operative Einsparungen",
    "Kurzere Wege, weniger Personal, weniger Energiekosten",
    "Keine Abhangigkeit von Phase B",
    "Halle 3 wird strategisches Asset",
    "Budget: 262\u2013460k \u20ac",
    "Umsetzung: Q1 2026 \u2013 Q1 2027",
]
set_ph_bullets(s9, 1, items_a, font_size=12, spacing=Pt(5))

# Right: Phase B
items_b = [
    "PHASE B \u2014 ZUKUNFT (SEPARAT)",
    "",
    "Nicht Teil dieses Budgets",
    "Keine Option fur Gundernhausen bestatigt",
    "4 Innovationspfade identifiziert",
    "CO2-RA: hochster TRL (6-7), Ellzee-Erfahrung",
    "PU-Recycling, LHS-TPC, FR-PU: neue Geschaftsfelder",
    "Halle 3 bleibt verfugbar \u2014 Entscheidung spater",
]
set_ph_bullets(s9, 2, items_b, font_size=12, spacing=Pt(5))

# ════════════════════════════════════════════════════════
# SLIDE 10: Next Steps
# ════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(L_CONTENT)
set_ph(s10, 12, "Empfehlung", font_size=11, color=GRAY66)
set_ph(s10, 0, "Nachste Schritte", font_size=24, bold=True)
set_ph(s10, 11, "NVH 2030 | Werk Gundernhausen | 26.02.2026 | Autoneum", font_size=8, color=GRAY_MED)

actions = [
    "1.  Vorschlag 4 als Basislayout formal bestatigen                          Werkleitung          Marz 2026",
    "2.  Budgetfreigabe Phase A (262\u2013460k \u20ac)                                       Controlling            Marz 2026",
    "3.  Infrastruktur-Bestandsaufnahme Halle 2                                   Technik                 April 2026",
    "4.  Ausschreibung Maschinenabbau (5 Maschinen)                           Einkauf                 Mai 2026",
    "5.  Sommerpause nutzen fur kritische Umbauten                              Produktion            Jul/Aug 2026",
    "",
    "Phase A rechnet sich unabhangig von Phase B allein durch operative Einsparungen.",
    "Halle 3 bleibt als strategisches Asset verfugbar \u2014 welche Option dort realisiert wird, kann spater entschieden werden.",
]
set_ph_bullets(s10, 1, actions, font_size=12, spacing=Pt(6))

# ── Save ──
prs.save(OUTPUT)
print(f"PPTX saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
